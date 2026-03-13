"""
Documents MCP Server

Provides Model Context Protocol tools to inspect and read local documents
(PDF, Excel, CSV, JSON, TXT, Markdown). Designed to be launched via
VS Code's MCP configuration.

No external LLM API keys required — the host AI (e.g. GitHub Copilot)
performs all analysis on the returned content, including images.
"""
import json
import os
import base64
from pathlib import Path
from typing import Any, Dict, List, Optional

import pandas as pd
import pdfplumber
from mcp.server.fastmcp import FastMCP

try:
    import openpyxl
    from openpyxl.drawing.image import Image as XLImage
except ImportError:  # pragma: no cover
    openpyxl = None
    XLImage = None

try:
    import docx  # optional dependency
except ImportError:  # pragma: no cover
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:  # optional OCR dependencies
    import pytesseract
    from PIL import Image
except ImportError:  # pragma: no cover
    pytesseract = None
    Image = None

# ----------------------------------------------------------------------------
# Configuration
# ----------------------------------------------------------------------------
DEFAULT_ROOT = Path(os.getenv("DOCUMENTS_ROOT", Path(__file__).parent))
ALLOWED_EXTENSIONS = {
    ".pdf",
    ".txt",
    ".md",
    ".json",
    ".csv",
    ".xlsx",
    ".xls",
    ".tsv",
    ".docx",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".bmp",
    ".tif",
    ".tiff",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}

server = FastMCP("document-analyzer")


def _resolve_path(file_path: str) -> Path:
    path = (DEFAULT_ROOT / file_path).resolve() if not Path(file_path).is_absolute() else Path(file_path).resolve()
    if not str(path).startswith(str(DEFAULT_ROOT.resolve())):
        raise ValueError("Access outside of root directory is not allowed")
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")
    if path.suffix.lower() not in ALLOWED_EXTENSIONS:
        raise ValueError(f"Unsupported file type: {path.suffix}")
    return path


def _read_text_file(path: Path, max_chars: int) -> str:
    data = path.read_text(encoding="utf-8", errors="ignore")
    return data[:max_chars]


def _read_json_file(path: Path, max_chars: int) -> str:
    data = json.loads(path.read_text(encoding="utf-8"))
    pretty = json.dumps(data, indent=2)
    return pretty[:max_chars]


def _read_csv_file(path: Path, max_rows: int, sep: str = ",") -> str:
    df = pd.read_csv(path, sep=sep, nrows=max_rows)
    return df.to_string(index=False)


def _read_excel_file(path: Path, sheet: Optional[str], start_row: int, num_rows: int) -> str:
    """Read Excel file with full openpyxl support for chart/image detection."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError("openpyxl is required for Excel files; install it with: pip install openpyxl")
    
    # Load workbook to inspect structure
    wb = load_workbook(path, data_only=True)
    available_sheets = wb.sheetnames
    
    # Validate requested sheet
    if sheet:
        if sheet not in available_sheets:
            wb.close()
            raise ValueError(
                f"Sheet '{sheet}' not found. Available sheets: {', '.join(available_sheets)}"
            )
        sheet_name = sheet
    else:
        sheet_name = available_sheets[0]
    
    ws = wb[sheet_name]
    
    # Detect charts and images (filter small logos)
    has_charts = len(ws._charts) > 0 if hasattr(ws, '_charts') else False
    
    # Count only significant images (>= 50x50 pixels)
    significant_image_count = 0
    if hasattr(ws, '_images'):
        for img in ws._images:
            # Check image dimensions
            if hasattr(img, 'width') and hasattr(img, 'height'):
                # openpyxl stores dimensions in EMUs (English Metric Units)
                # 914400 EMUs = 1 inch, typical 96 DPI means ~9525 EMUs per pixel
                width_pixels = img.width / 9525 if img.width else 0
                height_pixels = img.height / 9525 if img.height else 0
                # Only count images >= 200x200 pixels (filter logos)
                if width_pixels >= 200 and height_pixels >= 200:
                    significant_image_count += 1
            else:
                # If no dimensions, count it
                significant_image_count += 1
    
    has_images = significant_image_count > 0
    chart_count = len(ws._charts) if has_charts else 0
    image_count = significant_image_count
    
    # Try to read data
    try:
        df = pd.read_excel(path, sheet_name=sheet_name)
        if df.empty or df.dropna(how='all').empty:
            # Empty data sheet, likely chart-only
            wb.close()
            sheets_list = ', '.join(f"'{s}'" for s in available_sheets)
            header = f"[Available sheets: {sheets_list}]\nReading sheet: '{sheet_name}'\n\n"
            msg = (
                f"⚠️ VISUAL CONTENT DETECTED: Sheet contains {chart_count} chart(s) and {image_count} image(s) but no data table ⚠️\n\n"
                f"RECOMMENDATION: Convert Excel to PDF first, then use:\n"
                f"  visual_evaluate_document(file_path='<pdf_path>', enable_vision=true, sheet='{sheet_name}')\n"
                "to extract chart content with GPT-4 Vision.\n"
            )
            return header + msg
        
        subset = df.iloc[start_row : start_row + num_rows]
        wb.close()
        
        # Include metadata in output
        sheets_list = ', '.join(f"'{s}'" for s in available_sheets)
        header = f"[Available sheets: {sheets_list}]\nReading sheet: '{sheet_name}'\n"
        if has_charts or has_images:
            header += (
                f"\n⚠️ VISUAL CONTENT DETECTED: {chart_count} chart(s) and {image_count} image(s) in this sheet ⚠️\n"
                f"NOTE: Text/data extraction successful (see below), but charts/images not extracted.\n"
                f"RECOMMENDATION: To extract chart content, convert to PDF and use:\n"
                f"  visual_evaluate_document(file_path='<pdf_path>', enable_vision=true)\n"
            )
        header += "\n"
        return header + subset.to_string(index=False)
    
    except Exception as e:
        # Sheet might be chart-only or corrupted
        wb.close()
        sheets_list = ', '.join(f"'{s}'" for s in available_sheets)
        header = f"[Available sheets: {sheets_list}]\nReading sheet: '{sheet_name}'\n\n"
        msg = (
            f"⚠️ VISUAL CONTENT DETECTED: {chart_count} chart(s) and {image_count} image(s) ⚠️\n"
            f"[Cannot read sheet as data table: {str(e)}]\n\n"
            f"RECOMMENDATION: Convert Excel to PDF first, then use:\n"
            f"  visual_evaluate_document(file_path='<pdf_path>', enable_vision=true)\n"
            "to analyze this chart-only sheet with GPT-4 Vision.\n"
        )
        return header + msg


def _read_pdf_file(path: Path, page: int, max_pages: int, max_chars: int) -> str:
    """Read PDF with detection of images/charts to recommend vision analysis."""
    texts: List[str] = []
    visual_content_detected = []
    
    with pdfplumber.open(path) as pdf:
        total_pages = len(pdf.pages)
        start = max(page - 1, 0)
        end = min(start + max_pages, total_pages)
        
        for idx in range(start, end):
            page_obj = pdf.pages[idx]
            
            # Detect visual content on this page
            # Filter out small images (logos, icons) - only count images >= 200x200 pixels
            significant_images = []
            if hasattr(page_obj, 'images'):
                for img in page_obj.images:
                    # Check image dimensions
                    width = img.get('width', 0)
                    height = img.get('height', 0)
                    # Only count images that are large enough to contain data (>= 200x200)
                    # Logos and small icons are typically < 200 pixels
                    if width >= 200 and height >= 200:
                        significant_images.append(img)
            
            has_images = len(significant_images) > 0
            has_curves = len(page_obj.curves) > 0 if hasattr(page_obj, 'curves') else False
            
            # Extract text
            extracted = page_obj.extract_text() or ""
            
            # Check if page is mostly images (scanned or has charts)
            if has_images or (has_curves and len(extracted.strip()) < 100):
                image_count = len(significant_images)
                visual_content_detected.append(f"Page {idx + 1}: {image_count} image(s), possible chart/diagram")
            
            header = f"\n--- Page {idx + 1}/{total_pages} ---\n"
            texts.append(header + extracted)
    
    joined = "\n".join(texts)
    
    # Add recommendation if visual content detected
    if visual_content_detected:
        recommendation = (
            "\n\n⚠️ VISUAL CONTENT DETECTED ⚠️\n"
            + "\n".join(visual_content_detected)
            + f"\n\nRECOMMENDATION: Use visual_evaluate_document(file_path='{path.name}', enable_vision=true, page={page}) "
            "to extract charts, diagrams, and image content with GPT-4 Vision.\n"
        )
        joined = recommendation + joined
    
    return joined[:max_chars + 500]  # Allow extra space for recommendation


def _format_table(rows: List[List[str]]) -> str:
    if not rows:
        return ""
    num_cols = max(len(row) for row in rows)
    widths = [0] * num_cols
    for row in rows:
        for idx in range(num_cols):
            cell = row[idx] if idx < len(row) else ""
            widths[idx] = max(widths[idx], len(cell))
    lines: List[str] = []
    for row in rows:
        cells: List[str] = []
        for idx in range(num_cols):
            cell = row[idx] if idx < len(row) else ""
            cells.append(cell.ljust(widths[idx]))
        lines.append(" | ".join(cells).rstrip())
    return "\n".join(lines)


def _read_docx_file(path: Path, max_chars: int, include_tables: bool) -> str:
    """Read Word document with detection of images and embedded objects."""
    if docx is None:
        raise ImportError("python-docx is required to read DOCX files")
    
    document = docx.Document(str(path))
    
    # Detect images and embedded objects (filter out small logos)
    from PIL import Image as PILImage
    import io
    
    image_count = 0
    for rel in document.part.rels.values():
        if "image" in rel.target_ref:
            try:
                # Try to get image dimensions
                image_part = rel.target_part
                img_stream = io.BytesIO(image_part.blob)
                img = PILImage.open(img_stream)
                width, height = img.size
                # Only count images >= 200x200 pixels (filter logos and small icons)
                if width >= 200 and height >= 200:
                    image_count += 1
            except:
                # If we can't read dimensions, count it anyway
                image_count += 1
    
    # Count tables
    table_count = len(document.tables)
    
    # Extract text
    paragraphs = [p.text for p in document.paragraphs]
    text = "\n".join(paragraphs)
    
    # Add metadata header with visual content detection
    header_parts = []
    if image_count > 0:
        header_parts.append(f"{image_count} image(s)")
    if table_count > 0:
        header_parts.append(f"{table_count} table(s)")
    
    if header_parts:
        header = f"[Document contains: {', '.join(header_parts)}]\n"
        if image_count > 0:
            header += (
                f"\n⚠️ VISUAL CONTENT DETECTED: {image_count} image(s) in this Word document ⚠️\n"
                f"RECOMMENDATION: Convert to PDF first, then use:\n"
                f"  visual_evaluate_document(file_path='<pdf_path>', enable_vision=true)\n"
                "to extract content from images with GPT-4 Vision.\n"
            )
        header += "\n"
        text = header + text
    
    # Extract tables if requested
    if include_tables and document.tables:
        table_chunks: List[str] = []
        for idx, table in enumerate(document.tables, start=1):
            tables_rows: List[List[str]] = []
            for row in table.rows:
                tables_rows.append([cell.text.strip() for cell in row.cells])
            formatted = _format_table(tables_rows)
            table_chunks.append(f"Table {idx}:\n{formatted}")
        if table_chunks:
            text += "\n\n" + "\n\n".join(table_chunks)
    
    return text[:max_chars]


def _read_pptx_file(path: Path, max_chars: int) -> str:
    """Read PowerPoint file with detection of visuals."""
    if pptx is None:
        raise ImportError("python-pptx is required to read PPTX files")
        
    prs = pptx.Presentation(str(path))
    texts = []
    
    # Metadata
    slide_count = len(prs.slides)
    header = f"[PowerPoint Presentation: {slide_count} slides]\n"
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        title = slide.shapes.title.text if slide.shapes.title else "No Title"
        slide_text.append(f"--- Slide {i+1}: {title} ---")
        
        # Shapes text
        for shape in slide.shapes:
             if hasattr(shape, "text") and shape.text:
                 text_clean = shape.text.strip()
                 if text_clean and text_clean != title:
                     slide_text.append(text_clean)
             
             if shape.has_table:
                 table = shape.table
                 rows_data = []
                 for row in table.rows:
                     rows_data.append([cell.text_frame.text.strip() for cell in row.cells])
                 formatted = _format_table(rows_data)
                 slide_text.append(f"\n[Table]\n{formatted}\n")
                 
        texts.append("\n".join(slide_text))
        
    return header + "\n\n".join(texts)[:max_chars]


def _require_ocr_dependencies() -> None:
    if pytesseract is None or Image is None:
        raise ImportError("pytesseract and pillow are required for OCR; install them to enable this feature")



def _ocr_image(path: Path, lang: str) -> str:
    _require_ocr_dependencies()
    with Image.open(path) as img:
        return pytesseract.image_to_string(img, lang=lang)


def _ocr_pdf_pages(path: Path, page: int, max_pages: int, lang: str) -> str:
    _require_ocr_dependencies()
    texts: List[str] = []
    with pdfplumber.open(path) as pdf:
        total_pages = len(pdf.pages)
        start = max(page - 1, 0)
        end = min(start + max_pages, total_pages)
        for idx in range(start, end):
            page_obj = pdf.pages[idx]
            page_image = page_obj.to_image(resolution=200).original
            ocr_text = pytesseract.image_to_string(page_image, lang=lang)
            texts.append(f"--- OCR Page {idx + 1}/{total_pages} ---\n{ocr_text.strip()}")
    return "\n\n".join(texts)


def _encode_image_base64(image_path: Path) -> str:
    """Encode image file to base64 string."""
    with open(image_path, "rb") as img_file:
        return base64.b64encode(img_file.read()).decode('utf-8')


def _pdf_page_to_base64(pdf_path: Path, page_num: int) -> str:
    """Convert a PDF page to base64-encoded image."""
    _require_ocr_dependencies()
    with pdfplumber.open(pdf_path) as pdf:
        if page_num < 1 or page_num > len(pdf.pages):
            raise ValueError(f"Page {page_num} out of range (1-{len(pdf.pages)})")
        page = pdf.pages[page_num - 1]
        page_image = page.to_image(resolution=200).original
        from io import BytesIO
        buffer = BytesIO()
        page_image.save(buffer, format="PNG")
        return base64.b64encode(buffer.getvalue()).decode('utf-8')


def _get_image_mime(path: Path) -> str:
    """Return MIME type for an image file."""
    ext = path.suffix.lower()
    return {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".bmp": "image/bmp",
        ".tif": "image/tiff",
        ".tiff": "image/tiff",
    }.get(ext, "image/png")


@server.tool()
def list_documents(
    root: Optional[str] = None,
    pattern: str = "**/*",
    limit: int = 100,
) -> List[Dict[str, Any]]:
    """List documents under the configured root directory."""
    base = Path(root).resolve() if root else DEFAULT_ROOT.resolve()
    results: List[Dict[str, Any]] = []
    for path in base.glob(pattern):
        if path.is_file() and path.suffix.lower() in ALLOWED_EXTENSIONS:
            results.append({
                "path": str(path),
                "size_bytes": path.stat().st_size,
                "suffix": path.suffix,
            })
        if len(results) >= limit:
            break
    return results


@server.tool()
def document_info(file_path: str) -> Dict[str, Any]:
    """Return metadata about a specific document."""
    path = _resolve_path(file_path)
    info = path.stat()
    response: Dict[str, Any] = {
        "path": str(path),
        "size_bytes": info.st_size,
        "modified": info.st_mtime,
        "extension": path.suffix,
    }
    if path.suffix.lower() in {".xlsx", ".xls"}:
        try:
            sheets = pd.ExcelFile(path).sheet_names
        except Exception as exc:  # pragma: no cover
            sheets = [f"Error: {exc}"]
        response["sheets"] = sheets
    return response


@server.tool()
def read_document(
    file_path: str,
    max_chars: int = 8000,
    page: int = 1,
    max_pages: int = 2,
    sheet: Optional[str] = None,
    start_row: int = 0,
    num_rows: int = 50,
    include_tables: bool = False,
) -> str:
    """Read a portion of a document (PDF, Excel, CSV, JSON, TXT, DOCX).
    
    IMPORTANT: When presenting results to users, ALWAYS mention:
    - The folder/directory path (e.g., 'Well 1/Well Test/')
    - The complete document name
    This helps users understand which well and document type is being referenced.
    """
    path = _resolve_path(file_path)
    suffix = path.suffix.lower()
    
    # Add context header with folder and filename
    folder_path = str(path.parent.relative_to(DEFAULT_ROOT)) if path.is_relative_to(DEFAULT_ROOT) else str(path.parent)
    context_header = f"📁 Document: {folder_path}/{path.name}\n" + "="*80 + "\n\n"

    if suffix in {".txt", ".md"}:
        return context_header + _read_text_file(path, max_chars)
    if suffix == ".json":
        return context_header + _read_json_file(path, max_chars)
    if suffix in {".csv", ".tsv"}:
        sep = "\t" if suffix == ".tsv" else ","
        return context_header + _read_csv_file(path, num_rows, sep)
    if suffix in {".xlsx", ".xls"}:
        return context_header + _read_excel_file(path, sheet, start_row, num_rows)
    if suffix == ".pdf":
        return context_header + _read_pdf_file(path, page, max_pages, max_chars)
    if suffix == ".docx":
        return context_header + _read_docx_file(path, max_chars, include_tables)
    if suffix == ".pptx":
        return context_header + _read_pptx_file(path, max_chars)

    raise ValueError(f"Unsupported file type: {suffix}")


@server.tool()
def visual_evaluate_document(
    file_path: str,
    enable_ocr: bool = False,
    page: int = 1,
    max_pages: int = 1,
    ocr_lang: str = "eng",
) -> list:
    """Extract visual content from a document (PDF pages or images) and return
    it as inline images so the host AI can analyse charts, tables, and diagrams
    directly.  Optionally run local Tesseract OCR as well.

    Returns a list of content blocks (text and/or base64 images) that the
    calling AI model can interpret.
    """
    from mcp.types import TextContent, ImageContent

    path = _resolve_path(file_path)
    suffix = path.suffix.lower()

    folder_path = (
        str(path.parent.relative_to(DEFAULT_ROOT))
        if path.is_relative_to(DEFAULT_ROOT)
        else str(path.parent)
    )

    contents: list = []
    contents.append(
        TextContent(
            type="text",
            text=f"Document: {folder_path}/{path.name}",
        )
    )

    # ---- Return images so the host AI can see them ----
    if suffix in IMAGE_EXTENSIONS:
        mime = _get_image_mime(path)
        b64 = _encode_image_base64(path)
        contents.append(
            ImageContent(type="image", data=b64, mimeType=mime)
        )
    elif suffix == ".pdf":
        try:
            with pdfplumber.open(path) as pdf:
                total = len(pdf.pages)
                start = max(page - 1, 0)
                end = min(start + max_pages, total)
                for idx in range(start, end):
                    page_b64 = _pdf_page_to_base64(path, idx + 1)
                    contents.append(
                        TextContent(type="text", text=f"--- Page {idx + 1}/{total} ---")
                    )
                    contents.append(
                        ImageContent(type="image", data=page_b64, mimeType="image/png")
                    )
        except Exception as exc:
            contents.append(
                TextContent(type="text", text=f"Failed to render PDF pages: {exc}")
            )

        # Also include pdfplumber extracted text/tables as supplementary data
        try:
            with pdfplumber.open(path) as pdf:
                total = len(pdf.pages)
                start = max(page - 1, 0)
                end = min(start + max_pages, total)
                for idx in range(start, end):
                    pg = pdf.pages[idx]
                    text = pg.extract_text() or ""
                    tables = pg.extract_tables()
                    parts = [f"Page {idx + 1} extracted text:\n{text}"]
                    for ti, table in enumerate(tables, 1):
                        rows = "\n".join(
                            " | ".join(str(c or "") for c in row) for row in table
                        )
                        parts.append(f"Table {ti}:\n{rows}")
                    contents.append(
                        TextContent(type="text", text="\n\n".join(parts))
                    )
        except Exception:
            pass
    else:
        contents.append(
            TextContent(
                type="text",
                text=f"Visual extraction not supported for {suffix}; use read_document instead.",
            )
        )

    # ---- Optional local OCR ----
    if enable_ocr:
        try:
            if suffix == ".pdf":
                ocr_text = _ocr_pdf_pages(path, page, max_pages, ocr_lang)
            elif suffix in IMAGE_EXTENSIONS:
                ocr_text = _ocr_image(path, ocr_lang)
            else:
                ocr_text = f"OCR not supported for {suffix}"
            contents.append(
                TextContent(type="text", text=f"OCR result:\n{ocr_text}")
            )
        except Exception as exc:
            contents.append(
                TextContent(type="text", text=f"OCR failed: {exc}")
            )

    return contents


if __name__ == "__main__":
    server.run()
